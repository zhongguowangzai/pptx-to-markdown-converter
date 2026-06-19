"""
PPTX → MD 深度转换器 v3
改进：层级大纲检测、矩阵表格仅在均匀列时启用、保留顶部标题
"""
import sys, io, re, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Emu
from collections import defaultdict

def get_all_text(shape):
    texts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                texts.append({'text': t, 'level': para.level})
    return texts

class SlideParser:
    def __init__(self, slide):
        self.slide = slide
        self.title = ''
        self.subtitle = ''
        self.body = []
        self.sidebar = []
        self.footer = []
        self.tables = []
        self.flowchart_boxes = []
        self.all_text_shapes = []
        self.has_images = False
        self._top_title = ''  # 页面顶部的标题性文本

    def parse(self):
        for shape in self.slide.shapes:
            left_in = Emu(shape.left).inches
            top_in = Emu(shape.top).inches
            w_in = Emu(shape.width).inches
            h_in = Emu(shape.height).inches
            texts = get_all_text(shape)
            shape_type = shape.shape_type

            if shape_type == 13:  # PICTURE
                self.has_images = True
                continue
            if shape_type in (8, 9, 10, 11, 12) and not texts:
                continue
            if left_in > 12.0 and w_in < 2:  # 侧边栏
                for t in texts:
                    self.sidebar.append(t['text'])
                continue
            if top_in > 6.0 and h_in < 1.5:  # 底部
                for t in texts:
                    self.footer.append(t['text'])
                continue
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                if texts:
                    self.title = texts[0]['text']
                continue
            if shape.is_placeholder and shape.placeholder_format.idx in (1, 2):
                if shape.placeholder_format.idx == 2 and texts:
                    self.subtitle = texts[0]['text']
                    for t in texts[1:]:
                        self.body.append(t)
                else:
                    for t in texts:
                        self.body.append(t)
                continue
            if shape.has_table:
                table = shape.table
                rows_data = []
                for row in table.rows:
                    rows_data.append([cell.text.strip() for cell in row.cells])
                self.tables.append(rows_data)
                continue

            if texts and h_in > 0.2 and w_in > 0.3:
                # 保存顶部标题性文本（y < 0.5, 宽度中等的短文本）
                if top_in < 0.5 and 3 < w_in < 10 and len(texts) == 1 and len(texts[0]['text']) < 50:
                    self._top_title = texts[0]['text']
                    continue
                self.all_text_shapes.append({
                    'left': left_in, 'top': top_in, 'w': w_in, 'h': h_in,
                    'texts': texts,
                })

        self._process_text_shapes()

    def _check_uniform_columns(self, bands_data):
        """检查所有行是否列数一致，适合做表格"""
        col_counts = [len(row) for row in bands_data]
        if not col_counts:
            return False
        return len(set(col_counts)) == 1 and col_counts[0] >= 2

    def _detect_outline_structure(self):
        """检测层级大纲结构：左侧短标签 + 右侧长描述"""
        total = len(self.all_text_shapes)
        if total < 5:
            return False
        
        left_labels = 0   # left < 4, short text (<25 chars)
        right_descs = 0   # left > 5, long text (>30 chars)
        
        for s in self.all_text_shapes:
            all_text = ' '.join(t['text'] for t in s['texts'])
            if s['left'] < 4.0 and len(all_text) < 25:
                left_labels += 1
            if s['left'] > 5.0 and len(all_text) > 30:
                right_descs += 1
        
        return left_labels >= 3 and right_descs >= 2 and left_labels >= total * 0.2

    def _build_outline(self):
        """上下文感知：自上而下扫描，追踪当前所属层级；同 Y 带内标签优先"""
        # 按 0.5in Y 带分组，带内按 left 排序
        bands = defaultdict(list)
        for s in self.all_text_shapes:
            y_key = round(s['top'] * 2) / 2
            bands[y_key].append(s)
        
        output_lines = []
        context = {'L1': None, 'L2': None, 'L3': None}
        
        for y_key in sorted(bands.keys()):
            band_shapes = sorted(bands[y_key], key=lambda s: s['left'])
            band_lines = []
            
            for s in band_shapes:
                text = '\n'.join(t['text'] for t in s['texts'])
                left = s['left']
                
                # 跳过极短空文本
                if not text.strip():
                    continue
                
                if left < 1.5 and len(text) < 20:
                    context['L1'] = text
                    context['L2'] = None
                    context['L3'] = None
                    band_lines.append(f"\n**{text}**")
                    
                elif 1.5 <= left < 3.0 and len(text) < 15:
                    context['L2'] = text
                    context['L3'] = None
                    band_lines.append(f"- **{text}**")
                    
                elif 3.0 <= left < 5.5 and len(text) < 15:
                    # 三级标题（总额法/净额法/以前发生/未来发生），放宽到 5.5
                    context['L3'] = text
                    band_lines.append(f"    - **{text}**")
                    
                elif left >= 4.5 and len(text) > 10:
                    # 正文描述
                    if context['L3']:
                        band_lines.append(f"        {text}")
                    elif context['L2']:
                        band_lines.append(f"    {text}")
                    else:
                        band_lines.append(f"{text}")
                elif left >= 5.5 and len(text) < 15:
                    # 右侧额外的短标签（如位于 5.6 的"总额法："），作为 L3 处理
                    context['L3'] = text
                    band_lines.append(f"    - **{text}**")
                else:
                    band_lines.append(f"  {text}")
            
            output_lines.extend(band_lines)
        
        return output_lines

    def _process_text_shapes(self):
        if not self.all_text_shapes:
            return

        # 先检测是否为层级大纲结构
        if self._detect_outline_structure():
            outline_lines = self._build_outline()
            for line in outline_lines:
                if line:
                    self.body.append({'text': line, 'level': 0})
            return

        # 再检测是否为均匀矩阵表格
        bands = defaultdict(list)
        for s in self.all_text_shapes:
            if s['w'] >= 6 and len(s['texts']) == 1 and len(s['texts'][0]['text']) > 50:
                # 宽文本转义：极宽的单个长文本=内容说明，不入表格
                self.body.append(s['texts'][0])
                continue
            y_key = round(s['top'] * 2) / 2
            bands[y_key].append(s)

        sorted_bands = sorted(bands.items())
        bands_data = []
        for y, band_shapes in sorted_bands:
            sorted_band = sorted(band_shapes, key=lambda s: s['left'])
            row_texts = ['\n'.join(t['text'] for t in s['texts']) for s in sorted_band]
            if row_texts:
                bands_data.append(row_texts)

        if len(bands_data) >= 2 and self._check_uniform_columns(bands_data):
            self.tables.append(bands_data)
            if self._top_title:
                self.body.insert(0, {'text': self._top_title, 'level': 0})
            return

        # 都不是：流畅流程图检测
        if self._try_flowchart():
            return

        # 按Y排序输出
        shapes_sorted = sorted(self.all_text_shapes, key=lambda s: (round(s['top'], 1), s['left']))
        for s in shapes_sorted:
            for t in s['texts']:
                self.body.append(t)
        if self._top_title:
            self.body.insert(0, {'text': self._top_title, 'level': 0})

    def _try_flowchart(self):
        candidates = []
        for s in self.all_text_shapes:
            all_text = ' '.join(t['text'] for t in s['texts'])
            if len(all_text) <= 20 and len(s['texts']) == 1:
                if 1.5 < s['top'] < 6.5 and s['w'] >= 0.5 and s['h'] >= 0.3:
                    candidates.append(s)

        if len(candidates) < 3:
            return False

        candidates.sort(key=lambda s: s['top'])
        top_band = [c for c in candidates if c['top'] - candidates[0]['top'] < 0.8]
        if len(top_band) == 0 or len(top_band) >= len(candidates):
            return False

        self.flowchart_boxes = candidates
        return True


def build_flowchart_mermaid(boxes):
    if len(boxes) < 3:
        return None
    boxes_sorted = sorted(boxes, key=lambda b: b['top'])
    root_y = boxes_sorted[0]['top']
    roots, children = [], []
    for b in boxes_sorted:
        if b['top'] - root_y < 0.8:
            roots.append(b)
        else:
            children.append(b)
    if not children:
        return None

    mermaid = "```mermaid\ngraph TD\n"
    root_ids = []
    root_texts_set = set()
    for b in roots:
        text = b['texts'][0]['text'][:30].replace('\n', ' ').replace('"', "'")
        if not text or text in root_texts_set:
            continue
        root_texts_set.add(text)
        node_id = re.sub(r'[^\w\u4e00-\u9fff]', '_', text)[:30]
        mermaid += f'    {node_id}["{text}"]\n'
        root_ids.append(node_id)

    for b in children:
        text = b['texts'][0]['text'][:30].replace('\n', ' ').replace('"', "'")
        if not text:
            continue
        node_id = re.sub(r'[^\w\u4e00-\u9fff]', '_', text)[:30]
        if node_id in root_ids:
            node_id += '_child'
        mermaid += f'    {node_id}["{text}"]\n'
        for rid in root_ids:
            mermaid += f'    {rid} --> {node_id}\n'

    mermaid += "```\n"
    return mermaid


def build_md(prs):
    output = []
    for i, slide in enumerate(prs.slides, 1):
        parser = SlideParser(slide)
        parser.parse()

        output.append(f"\n<!-- Slide {i} -->\n")
        if parser.sidebar:
            output.append(f"> 📖 {' '.join(parser.sidebar)}\n\n")
        if parser.title:
            output.append(f"## {parser.title}\n")
        if parser.subtitle:
            output.append(f"### {parser.subtitle}\n")

        for t in parser.body:
            txt = t['text']
            lvl = t['level']
            # 跳过已经包含markdown格式的行
            if txt.startswith('**') or txt.startswith('  -') or txt.startswith('    '):
                output.append(f"{txt}\n")
            elif txt.endswith('？') or txt.endswith('?'):
                output.append(f"\n**{txt}**\n")
            elif lvl == 0:
                output.append(f"{txt}\n")
            else:
                output.append(f"{'  ' * lvl}- {txt}\n")

        for table in parser.tables:
            output.append("\n")
            for row_idx, row in enumerate(table):
                output.append('| ' + ' | '.join(row) + ' |\n')
                if row_idx == 0:
                    output.append('|' + '|'.join(['---'] * len(row)) + '|\n')
            output.append("\n")

        if parser.footer:
            output.append(f"\n> 💡 {' '.join(parser.footer)}\n")

        if parser.flowchart_boxes:
            fc = build_flowchart_mermaid(parser.flowchart_boxes)
            if fc:
                output.append(f"\n{fc}\n")

        output.append("\n---\n")

    return ''.join(output)


if __name__ == '__main__':
    PPTX = sys.argv[1] if len(sys.argv) > 1 else None
    OUT = sys.argv[2] if len(sys.argv) > 2 else (os.path.splitext(PPTX)[0] + '_精确版.md' if PPTX else None)
    if not PPTX or not os.path.exists(PPTX):
        print("用法: python pptx_convert_v3.py <input.pptx> [output.md]")
        sys.exit(1)

    prs = Presentation(PPTX)
    md = build_md(prs)
    with open(OUT, 'w', encoding='utf-8') as f:
        f.write(md)
    print(f"✅ 已生成: {OUT}")
    print(f"📄 共 {len(prs.slides)} 页")
