"""
最精确 PPT → MD 深度转换器
技术栈: python-pptx + placeholder.idx + 位置感知 + Mermaid流程图重建
用法: python pptx_deep_convert.py <input.pptx> [output.md]
"""
import sys, io, re, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Emu
from lxml import etree
from collections import defaultdict

PPTX = sys.argv[1] if len(sys.argv) > 1 else None
OUT = sys.argv[2] if len(sys.argv) > 2 else (os.path.splitext(PPTX)[0] + '_精确版.md' if PPTX else None)

if not PPTX or not os.path.exists(PPTX):
    print("用法: python pptx_deep_convert.py <input.pptx> [output.md]")
    sys.exit(1)

prs = Presentation(PPTX)

def get_all_text(shape):
    """提取形状所有文本"""
    texts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                texts.append({'text': t, 'level': para.level})
    return texts

def parse_slide(slide):
    """解析单页：返回 (title, subtitle, body_lines, sidebar_lines, footer_lines, flowchart_boxes, tables)"""
    title = ''
    subtitle = ''
    body = []
    sidebar = []   # 教材XX页
    footer = []    # 底部总结标注
    flowchart_boxes = []  # 流程图节点
    tables = []
    
    for shape in slide.shapes:
        left = Emu(shape.left).inches
        top = Emu(shape.top).inches
        w = Emu(shape.width).inches
        h = Emu(shape.height).inches
        
        texts = get_all_text(shape)
        xml = etree.tostring(shape._element, encoding='unicode')
        has_fill = 'solidFill' in xml or 'gradFill' in xml
        
        # 1. 标题占位符 (idx=0)
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            if texts:
                title = texts[0]['text']
            continue
        
        # 2. 侧边栏：右边缘 (>12in)
        if left > 12.0:
            for t in texts:
                sidebar.append(t['text'])
            continue
        
        # 3. 底部标注：y>6.0 且高度<1
        if top > 6.0 and h < 1.0:
            for t in texts:
                footer.append(t['text'])
            continue
        
        # 4. 表格
        if shape.has_table:
            table = shape.table
            rows_data = []
            for row in table.rows:
                rows_data.append([cell.text.strip() for cell in row.cells])
            tables.append(rows_data)
            continue
        
        # 5. 表格占位符 (idx=4)
        if shape.is_placeholder and shape.placeholder_format.idx == 4:
            # 可能是表格占位符
            table_xml = etree.tostring(shape._element, encoding='unicode')
            if 'graphicFrame' in table_xml:
                continue
        
        # 6. 子标题/正文
        if texts:
            # 正文占位符 (idx=1)
            if shape.is_placeholder and shape.placeholder_format.idx == 1:
                for t in texts:
                    body.append(t)
                continue
            
            # 子标题位置 (idx=2)
            if shape.is_placeholder and shape.placeholder_format.idx == 2:
                subtitle = texts[0]['text']
                for t in texts[1:]:
                    body.append(t)
                continue
            
            # 流程图检测：单一短文本(<=15字符,1段) + 非占位符 + 内容区
            is_connector = shape.shape_type == 8
            is_decor = (not texts) and w < 0.4
            short_single = texts and len(texts[0]['text']) <= 15 and len(texts) == 1
            
            if (not is_connector and not is_decor and not shape.is_placeholder 
                and short_single and top > 1.5 and top < 6.5 and w >= 0.5 and h >= 0.3):
                flowchart_boxes.append({
                    'text': texts[0]['text'] if texts else '',
                    'left': left, 'top': top, 'width': w, 'height': h,
                })
                continue
            
            # 其余文本归入正文
            for t in texts:
                body.append(t)
        
        # 7. 装饰性连接线（忽略）
    
    return title, subtitle, body, sidebar, footer, flowchart_boxes, tables


def build_flowchart_mermaid(boxes):
    """从流程图盒子构建 Mermaid。根节点=最上方, 子节点=下方全部"""
    if len(boxes) < 3:
        return None
    
    # 排序：按 Y 坐标
    boxes_sorted = sorted(boxes, key=lambda b: b['top'])
    
    # 找到根节点（最上方）和子节点（其余）
    # 如果第一个 box 的 Y 与第二个差距 >0.5in，则第一个是根节点
    root_y = boxes_sorted[0]['top']
    roots = []
    children = []
    
    for b in boxes_sorted:
        if b['top'] - root_y < 0.5:  # 0.5in 内算同一层
            roots.append(b)
        else:
            children.append(b)
    
    if not children:
        return None
    
    # 构建 Mermaid
    mermaid = "```mermaid\ngraph TD\n"
    
    # 根节点
    root_ids = []
    for b in roots:
        text = b['text'].replace('\n', ' ').replace('"', "'")
        if not text: continue
        node_id = re.sub(r'[^\w\u4e00-\u9fff]', '_', text)[:30]
        mermaid += f'    {node_id}["{text}"]\n'
        root_ids.append(node_id)
    
    # 子节点
    for b in children:
        text = b['text'].replace('\n', ' ').replace('"', "'")
        if not text: continue
        node_id = re.sub(r'[^\w\u4e00-\u9fff]', '_', text)[:30]
        mermaid += f'    {node_id}["{text}"]\n'
        for rid in root_ids:
            mermaid += f'    {rid} --> {node_id}\n'
    
    mermaid += "```\n"
    return mermaid


def build_md(prs):
    output = []
    
    for i, slide in enumerate(prs.slides, 1):
        title, subtitle, body, sidebar, footer, flowchart_boxes, tables = parse_slide(slide)
        
        output.append(f"\n<!-- Slide {i} -->\n")
        
        # 侧边栏
        if sidebar:
            output.append(f"> 📖 {' '.join(sidebar)}\n\n")
        
        # 标题
        if title:
            output.append(f"## {title}\n")
        if subtitle:
            output.append(f"### {subtitle}\n")
        
        # 正文
        for t in body:
            txt = t['text']
            lvl = t['level']
            if txt.endswith('？') or txt.endswith('?'):
                output.append(f"\n**{txt}**\n")
            elif lvl == 0:
                output.append(f"{txt}\n")
            else:
                output.append(f"{'  ' * lvl}- {txt}\n")
        
        # 表格
        for table in tables:
            output.append("\n")
            for row_idx, row in enumerate(table):
                output.append('| ' + ' | '.join(row) + ' |\n')
                if row_idx == 0:
                    output.append('|' + '|'.join(['---'] * len(row)) + '|\n')
            output.append("\n")
        
        # 页脚
        if footer:
            output.append(f"\n> 💡 {' '.join(footer)}\n")
        
        # 流程图
        if flowchart_boxes:
            fc = build_flowchart_mermaid(flowchart_boxes)
            if fc:
                output.append(f"\n{fc}\n")
        
        output.append("\n---\n")
    
    return ''.join(output)


md = build_md(prs)

with open(OUT, 'w', encoding='utf-8') as f:
    f.write(md)

print(f"✅ 已生成: {OUT}")
print(f"📄 共 {len(prs.slides)} 页")
